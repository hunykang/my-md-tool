import streamlit as st
from markitdown import MarkItDown
import os
import tempfile

st.set_page_config(page_title="문서 변환기", page_icon="📝")

st.title("📝 문서 ➡️ Markdown 변환기 - 쌍칼버전")
st.write("파워포인트(PPTX) 또는 엑셀(XLSX, XLS) 파일을 업로드하면 마크다운(.md) 파일로 변환해 줍니다.")

# 파일 업로더
uploaded_file = st.file_uploader("여기에 파일을 드래그하거나 선택하세요", type=["pptx", "xlsx", "xls"])

if uploaded_file is not None:
    # 파일을 markitdown이 읽을 수 있도록 임시 파일로 저장
    file_extension = uploaded_file.name.split('.')[-1]
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}") as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_file_path = tmp_file.name

    st.info(f"선택된 파일: **{uploaded_file.name}**")
    
    # 변환 시작 버튼
    if st.button("Markdown으로 변환하기"):
        with st.spinner('문서를 분석하고 변환하는 중입니다...'):
            try:
                # MarkItDown 인스턴스 생성 및 변환
                md = MarkItDown()
                result = md.convert(tmp_file_path)
                md_content = result.text_content
                
                if not md_content or md_content.strip() == "":
                    st.warning("변환이 완료되었으나 추출된 텍스트가 없습니다. 파일 내용이 비어있거나 이미지로만 구성되어 있을 수 있습니다.")
                else:
                    st.success("🎉 변환이 완료되었습니다!")
                    
                    base_name = os.path.splitext(uploaded_file.name)[0]
                    new_file_name = f"{base_name}.md"
                    
                    # 이미지가 포함된 PPT의 경우 ZIP으로 묶어서 제공
                    if file_extension.lower() == "pptx":
                        import pptx
                        import re
                        import io
                        import zipfile

                        image_files = {}
                        try:
                            # python-pptx를 사용해 원본 PPT 파일의 이미지 추출
                            presentation = pptx.Presentation(tmp_file_path)
                            for slide in presentation.slides:
                                for shape in slide.shapes:
                                    # PICTURE (13) 이거나 PLACEHOLDER (14) 중 이미지가 있는 경우
                                    if getattr(shape, "shape_type", None) == 13 or (getattr(shape, "shape_type", None) == 14 and hasattr(shape, "image")):
                                        # markitdown 모듈과 똑같은 방식으로 이름 지정
                                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                                        try:
                                            image_files[filename] = shape.image.blob
                                        except Exception:
                                            pass
                        except Exception as e:
                            st.warning(f"이미지 추출 중 문제가 발생했습니다: {e}")

                        if len(image_files) > 0:
                            st.info(f"💡 PPT에 포함된 총 **{len(image_files)}개**의 그림 파일을 찾았습니다. 아래 버튼을 눌러 마크다운과 그림 파일을 ZIP으로 한 번에 받으세요!")
                            
                            # ZIP 메모리 버퍼 생성
                            zip_buffer = io.BytesIO()
                            with zipfile.ZipFile(zip_buffer, "a", zipfile.ZIP_DEFLATED, False) as zip_file:
                                # 마크다운 파일 추가 (.encode('utf-8') 필수)
                                zip_file.writestr(new_file_name, md_content.encode('utf-8'))
                                # 이미지 파일들 추가
                                for img_name, img_blob in image_files.items():
                                    zip_file.writestr(img_name, img_blob)

                            st.download_button(
                                label="📦 Markdown + 그림 파일 다운로드 (ZIP)",
                                data=zip_buffer.getvalue(),
                                file_name=f"{base_name}_변환결과.zip",
                                mime="application/zip"
                            )
                        
                        # 텍스트만 다운로드 하는 기본 버튼
                        st.download_button(
                            label="⬇️ Markdown 단일 파일만 다운로드 (.md)",
                            data=md_content,
                            file_name=new_file_name,
                            mime="text/markdown"
                        )
                    else:
                        # 엑셀 등의 기타 파일
                        st.download_button(
                            label="⬇️ Markdown 파일 다운로드 (.md)",
                            data=md_content,
                            file_name=new_file_name,
                            mime="text/markdown"
                        )

                    # 미리보기 제공
                    with st.expander("👀 마크다운 내용 미리보기", expanded=True):
                        st.text(md_content)
                        
            except Exception as e:
                st.error(f"변환 중 오류가 발생했습니다: {str(e)}")
            finally:
                # 임시 파일 정리
                if os.path.exists(tmp_file_path):
                    os.unlink(tmp_file_path)
